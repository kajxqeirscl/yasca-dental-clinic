# -*- coding: utf-8 -*-
"""
Yasca Dental Clinic - Final Proje Sunumu Yeniden Duzenleme
Mevcut PPTX'i temel alarak speaker notes eklenmiş, 
icerik daha anlasilir hale getirilmis yeni sunum olusturur.
"""
import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn, nsdecls
from pptx.oxml import parse_xml
import copy
import os

# ─── Color Palette (from original) ───
DARK_BG      = RGBColor(0x0A, 0x1F, 0x3B)  # Slide background (deep navy)
MID_NAVY     = RGBColor(0x0C, 0x42, 0x70)  # Headings
ACCENT_BLUE  = RGBColor(0x0E, 0x51, 0x89)  # Numbers, accents
LIGHT_BLUE   = RGBColor(0x9E, 0xB8, 0xD2)  # Subtitle text
LIGHTER_BLUE = RGBColor(0xBF, 0xD2, 0xE6)  # Italic subtitles
CARD_BLUE_BG = RGBColor(0xCF, 0xDD, 0xEC)  # Light card titles
CREAM        = RGBColor(0xF3, 0xEF, 0xE0)  # Big titles on dark
WHITE_BLUE   = RGBColor(0xE8, 0xEF, 0xF6)  # Body text on dark cards
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT     = RGBColor(0x2B, 0x2B, 0x2B)  # Body text on light
GRAY_TEXT     = RGBColor(0x5A, 0x5A, 0x5A)  # Footnotes/italic

# ─── Dimensions ───
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)
MARGIN  = Emu(548640)     # ~0.6 inch
LINE_Y  = Emu(987552)
TITLE_Y = Emu(274320)
BODY_Y  = Emu(1417320)

# Source presentation (for images and base layout)
SRC_PATH = r"c:\Users\Ali\yasca-dental-clinic\Yasca_Final_Proje_Raporu - Kopya.pptx"
OUT_PATH = r"c:\Users\Ali\yasca-dental-clinic\Yasca_Final_Proje_Raporu_Duzenlenmis.pptx"

src_prs = Presentation(SRC_PATH)
# We'll copy the source to keep images intact, then modify
prs = Presentation(SRC_PATH)

# ─── Helper Functions ───

def set_slide_bg(slide, color=DARK_BG):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height):
    """Add a text box and return it."""
    return slide.shapes.add_textbox(left, top, width, height)

def set_text(tf, text, font_size=Pt(14), color=DARK_TEXT, bold=False, italic=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    """Set text in a text frame, clearing existing content."""
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    return p

def add_para(tf, text, font_size=Pt(14), color=DARK_TEXT, bold=False, italic=False, space_before=Pt(4), space_after=Pt(2), alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    """Add a new paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    return p

def add_bullet(tf, text, font_size=Pt(13), color=DARK_TEXT, bullet_char='•'):
    """Add a bullet point."""
    return add_para(tf, f"{bullet_char}  {text}", font_size=font_size, color=color)

def add_notes(slide, notes_text):
    """Add speaker notes to a slide."""
    if not slide.has_notes_slide:
        slide.notes_slide  # This creates the notes slide
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.clear()
    # Split by lines and add
    for i, line in enumerate(notes_text.strip().split('\n')):
        if i == 0:
            tf.paragraphs[0].text = line.strip()
        else:
            p = tf.add_paragraph()
            p.text = line.strip()

def clear_shapes_except_bg(slide):
    """Remove all shapes except background rectangle."""
    shapes_to_remove = []
    for shape in slide.shapes:
        shapes_to_remove.append(shape)
    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)

# ─── Extract images from source slides ───
def extract_image_from_slide(src_slide, image_name_contains=None):
    """Extract image blob from source slide."""
    for shape in src_slide.shapes:
        if hasattr(shape, "image"):
            try:
                return shape.image.blob, shape.image.content_type
            except:
                pass
    return None, None

# Store original images from source
src_slides = list(src_prs.slides)
images = {}
for i, s in enumerate(src_slides):
    for shape in s.shapes:
        if hasattr(shape, "image"):
            try:
                blob = shape.image.blob
                ct = shape.image.content_type
                images[f"slide_{i+1}"] = (blob, ct)
            except:
                pass

# ═══════════════════════════════════════════════════════════════════
# Now we'll modify each slide of the copied presentation
# Adding speaker notes to EVERY slide
# ═══════════════════════════════════════════════════════════════════

slides = list(prs.slides)

# ─── SLIDE 1: Kapak ───
add_notes(slides[0], """
Merhaba hocam ve arkadaşlar. Biz Grup 4 olarak "Yaşca Diş Kliniği Yönetim Sistemi" projemizi sunacağız.

Projemiz, diş kliniklerinin günlük operasyonlarını dijitalleştirmek amacıyla geliştirdiğimiz çok-kiracılı (multi-tenant) bir SaaS uygulamasıdır.

Django REST Framework ve React.js kullanarak geliştirdik. Sistem, her kliniğin verisini PostgreSQL şema izolasyonu ile güvenli bir şekilde ayırıyor.

Bu sunumda, analiz aşamasından test sonuçlarına kadar tüm proje sürecini sizlerle paylaşacağız.
""")

# ─── SLIDE 2: Raporun Genel Akışı ───
add_notes(slides[1], """
Sunumumuz 8 ana bölümden oluşuyor:

1. Giriş ve Proje Kapsamı: Projenin ne olduğu, hedefleri ve hedef kullanıcıları.
2. Yazılım Gereksinimleri: 6 fonksiyonel ve 6 fonksiyonel olmayan ister.
3. UML Tasarımı: 6 sekans diyagramı ve 2 sınıf diyagramı ile tasarım modelleme.
4. Teknoloji Karar Analizi (DAR): Django, React, PostgreSQL seçimlerinin nedenleri ve karar matrisleri.
5. İç ve Dış Arayüzler: API kontratları, Swagger dokümantasyonu ve multi-tenant mimari.
6. Birim Test Sonuçları: 5 kritik sınıfta 54 birim testin tamamının PASSED olması.
7. Bakım ve Gelecek Önerileri: Refactoring, yeni özellikler ve versiyon planı.
8. Sonuç: Projenin özeti ve temel çıktılar.

Her bölümü sırasıyla ele alacağız.
""")

# ─── SLIDE 3: Giriş ve Proje Kapsamı ───
add_notes(slides[2], """
Projenin Amacı:
Diş klinikleri hâlâ birçoğu kağıt tabanlı veya kopuk sistemlerle çalışıyor. Bu durum veri kaybı, randevu çakışması ve güvenlik açıkları gibi sorunlara yol açıyor.

Yaşca Diş Kliniği Yönetim Sistemi, tüm bu sorunları çözmek için tasarlandı. Hasta kaydı, randevu yönetimi, tedavi takibi, dijital diş şeması ve finansal takip modüllerini tek bir platformda sunuyor.

En önemli özelliği: Çok-kiracılı (multi-tenant) SaaS mimarisi. Yani bir tane uygulama kurulur ama her klinik kendi ayrı alanında, kendi verileriyle çalışır. PostgreSQL şema izolasyonu sayesinde klinikler arası veri sızıntısı mimari düzeyde imkânsızdır.

Hedef Kullanıcılar:
- Yönetici (Admin): Klinik ayarları, personel yönetimi, tüm verilere erişim
- Hekim: Kendi hastaları, tedavi girişi, dijital odontogram, dashboard
- Asistan: Hasta kaydı, randevu oluşturma, ödeme girişi

Kritik Başarı Kriterleri olarak: Randevu çakışması sıfır olmalı, veri sızıntısı olmamalı, API yanıt süresi 2 saniyenin altında olmalı, en az 54 birim test geçmeli ve tüm endpointler Swagger ile dokümante edilmeli.
""")

# ─── SLIDE 4: Fonksiyonel İsterler ───
add_notes(slides[3], """
Proje Adımı 1'de belirlediğimiz 6 fonksiyonel ister şunlardır:

FR-01 - Kimlik Doğrulama: Kullanıcı adı ve parola ile sisteme giriş yapılır, JWT tabanlı access ve refresh token çifti döner. Her API çağrısında bu token kontrol edilir.

FR-02 - Hasta Yönetimi: Yeni hasta kaydı oluşturma, tıbbi özgeçmiş yani anamnez ekleme, güncelleme ve listeleme işlemleri. Ad, soyad ve telefon zorunlu alanlardır.

FR-03 - Randevu Yönetimi: Bu en kritik isterimiz. Randevu oluşturulurken sistem otomatik olarak çakışma kontrolü yapar. Aynı hekime aynı tarih ve saatte ikinci bir randevu girilmesini engeller.

FR-04 - Dashboard: Klinik personeli güne başlarken bugünkü randevuları, tamamlanan işlem sayısını ve aktif hasta sayısını tek bir ekrandan görür.

FR-05 - Tedavi ve Finans: Tedavi türü seçimi, diş numarası girişi, fiyat belirleme ve ödeme kaydı entegre çalışır.

FR-06 - Sistem Yönetimi: Sadece Admin rolündeki kullanıcı klinik çalışma saatlerini ve ayarlarını güncelleyebilir. Yetki kontrolü 403 Forbidden ile sağlanır.
""")

# ─── SLIDE 5: NFR ───
add_notes(slides[4], """
Fonksiyonel olmayan isterlerimiz, mimari kararlarımızın doğrudan gerekçesini oluşturur:

NFR-01 Performans: API sorgu süresinin 2 saniyenin altında olması hedeflenmiştir. Dashboard gibi aggregate veri dönen endpointlerde backend tarafında önceden hesaplama yapılır.

NFR-02 Güvenlik: JWT tabanlı kimlik doğrulama ve RBAC yani Role-Based Access Control ile yetkilendirme. Hekim sadece kendi hastalarını görebilir.

NFR-03 Veri İzolasyonu: django-tenants kütüphanesi ile her klinik PostgreSQL'de ayrı bir şemada tutulur. Bu fiziksel izolasyondur, yazılımsal değil.

NFR-04 Ölçeklenebilirlik: Birden fazla kliniğin eşzamanlı kullanımı desteklenir. Yeni bir klinik eklemek, yeni bir PostgreSQL şeması oluşturmak kadar basittir.

NFR-05 Bakım Yapılabilirlik: Modüler mimari sayesinde her modül bağımsız geliştirilebilir. Yüksek test kapsamı ile değişikliklerin yan etkilerini kontrol ediyoruz.

NFR-06 Kullanılabilirlik: Türkçe arayüz, sezgisel kullanım. React bileşen tabanlı mimarisi ile tutarlı kullanıcı deneyimi.
""")

# ─── SLIDE 6: UML Genel ───
add_notes(slides[5], """
Proje Adımı 1'de sistemi Boundary-Control-Entity ayrımıyla modelledik.

Sol tarafta Sekans Diyagramlarımız var - bunlar sistemin dinamik modelini gösterir:
- Her bir fonksiyonel gereksinim için ayrı bir sekans diyagramı çizdik, toplamda 6 tane.
- Bu diyagramlar nesneler arası mesaj akışlarını, hata durumlarını ve alternatif akışları modelliyor.

Sağ tarafta Sınıf Diyagramlarımız var - bunlar statik modeli gösterir:
- Domain Katmanı: 8 veri varlığı ve aralarındaki ilişkiler (Association, Composition, Aggregation)
- Uygulama Katmanı: Frontend bileşenleri (Boundary) ve Backend kontrolcüler (Control)

Tasarım tutarlılığı konusunda 3 önemli noktamız var:
1. Sekans diyagramlarındaki tüm mesajlar (fetchDashboardToday, createAppointment gibi), sınıf diyagramındaki ApiService metodları olarak birebir tanımlanmıştır.
2. Sorumlulukların ayrılığı: Boundary doğrudan Entity ile konuşmaz, araya Control girer. Bu MVC/MVT mimarisine uygundur.
3. Çokluk doğrulaması: Clinic ile Patient arasında 1'e çok ilişki, Patient ile Anamnesis arasında 1'e 1 ilişki gibi gerçek dünyayı yansıtan modelleme.
""")

# ─── SLIDE 7: Domain Sınıf Diyagramı (IMAGE) ───
add_notes(slides[6], """
Bu sınıf diyagramı, sistemin veritabanı tablolarının nesne tabanlı karşılıklarını gösterir.

Soldaki görsel, asıl UML diyagramımızdır. Sağ tarafta ilişki türlerini özetledik:

- Clinic → Patient: 1'e çok Association. Bir klinikte birden fazla hasta kayıtlıdır.
- Clinic → CustomUser: 1'e çok Association. Bir klinikte birden fazla personel bulunur.
- Patient → Anamnesis: 1'e 1 Composition. Her hastanın en fazla bir anamnez kaydı vardır, hasta silinirse anamnez de silinir.
- Patient → Appointment: 1'e çok Composition. Hasta silindiğinde randevuları da silinir.
- Appointment → Treatment: 1'e çok Aggregation. Bir randevuda birden fazla tedavi yapılabilir.
- Treatment → Payment: 1'e 1 Association. Her tedaviye bir ödeme kaydı bağlanır.
- TreatmentType → Treatment: 1'e çok. Bir tedavi türü birden fazla tedavide kullanılabilir.

Bu ilişkiler kodda ForeignKey, OneToOneField gibi Django ORM yapılarıyla birebir eşlenmiştir.
""")

# ─── SLIDE 8: Uygulama Katmanı Sınıf Diyagramı (IMAGE) ───
add_notes(slides[7], """
Bu diyagram, uygulama katmanını gösterir. 3 stereotip kullanıyoruz:

Boundary (Arayüz Bileşenleri): LoginPage, DashboardUI, PatientDialog, AppointmentDialog gibi React bileşenleri. Bunlar kullanıcının etkileştiği arayüzlerdir.

Control (İş Mantığı): İki katmanda çalışır:
- Frontend tarafında ApiService: Tüm HTTP çağrılarını merkezi olarak yönetir, JWT token'ı otomatik ekler.
- Backend tarafında ViewSet ve Serializer: Gelen istekleri karşılar, doğrulama yapar, iş kurallarını uygular.

Önemli tasarım kuralımız: Boundary asla doğrudan Entity ile konuşmaz. Her zaman araya Control katmanı girer. Bu sayede:
1. Güvenlik kontrolleri merkezi yapılır
2. İş kuralları tek noktada uygulanır
3. Arayüz değişiklikleri backend'i etkilemez

Bu MVC/MVT mimarisi prensiplerine uygun bir tasarımdır.
""")

# ─── SLIDE 9: Sekans SD-02 (IMAGE) ───
add_notes(slides[8], """
Bu sekans diyagramı, Hasta Kaydı ve Anamnez Güncelleme akışını gösterir (FR-02).

Akış şöyle ilerler:
1. Asistan veya Hekim, PatientDialog arayüzünden hasta bilgilerini girer.
2. PatientDialog, ApiService üzerinden POST /api/patients/ endpointine istek atar.
3. Backend'de PatientViewSet bu isteği karşılar ve PatientSerializer ile doğrulama yapar.
4. Ad, soyad ve telefon zorunlu alan kontrolü yapılır.
5. Hasta kaydı oluşturulur.
6. Eğer anamnez verisi de gönderilmişse, Anamnesis modeli zincirleme olarak oluşturulur.
7. Başarılı yanıt 201 Created ile döner.

Bu diyagramda dikkat edilmesi gereken: İç içe (nested) veri yapısı ile hasta ve anamnez tek istekte oluşturulabiliyor. PatientSerializer'ın create() metodunda bu nested write işlemi gerçekleştirilir.
""")

# ─── SLIDE 10: Sekans SD-03 (IMAGE) ───
add_notes(slides[9], """
Bu sekans diyagramı, projemizin en kritik akışını gösterir: Randevu Oluşturma ve Çakışma Kontrolü (FR-03).

Akış:
1. Asistan, AppointmentDialog'dan randevu bilgilerini girer: hasta, hekim, tarih ve saat.
2. ApiService, POST /api/appointments/ endpointine istek atar.
3. AppointmentViewSet isteği alır ve AppointmentCreateSerializer'a yönlendirir.
4. Serializer'ın validate() metodu devreye girer. Burada çakışma kontrolü yapılır:
   - Aynı hekime, aynı tarih ve saatte, durumu "Scheduled" olan başka bir randevu var mı?
   - Appointment.objects.filter(doctor=doctor, date=date, time=time, status='scheduled')
5. Eğer çakışma varsa: ValidationError fırlatılır, 400 Bad Request döner.
6. Eğer çakışma yoksa: Randevu "scheduled" statüsünde kaydedilir, 201 Created döner.

Diyagramda "alt" fragment ile bu iki alternatif akış gösterilmiştir. İptal edilmiş bir randevunun slotu bloke etmemesi de önemli bir detaydır - bunu testlerde de doğruladık.
""")

# ─── SLIDE 11: Sekans SD-04 (IMAGE) ───
add_notes(slides[10], """
Dashboard Günlük Özet akışı (FR-04).

Klinik personeli güne başlarken ana paneli açtığında:
1. DashboardUI bileşeni, ApiService üzerinden GET /api/dashboard/today/ endpointini çağırır.
2. Backend'de DashboardView, birden fazla sorguyu tek seferde çalıştırır:
   - Bugüne ait randevuları filtreler
   - Tamamlanan işlem sayısını hesaplar
   - Aktif hasta sayısını count eder
3. Bu aggregate veri tek bir JSON yanıtında döner.

Bu tasarım kararının nedeni performanstır. Frontend 3 ayrı API çağrısı yapmak yerine, backend tarafında hesaplanan özet veriyi tek seferde alır. Bu "Backend for Frontend" pattern'ine uygun bir yaklaşımdır.
""")

# ─── SLIDE 12: Sekans SD-06 (IMAGE) ───
add_notes(slides[11], """
Klinik Ayarları Güncelleme akışı (FR-06).

Bu akış, yetki kontrolünün (Authorization) nasıl çalıştığını gösterir:

1. Kullanıcı, ClinicSettingsUI ekranına erişir.
2. ApiService, PUT/PATCH /api/clinic-settings/ endpointine istek atar.
3. Backend'de ClinicSettingsViewSet isteği alır.
4. İlk olarak IsAdminUser permission sınıfı devreye girer.
5. Eğer kullanıcı Admin değilse: 403 Forbidden döner. İşlem yapılamaz.
6. Eğer kullanıcı Admin ise: Çalışma saatleri güncellenir, 200 OK döner.

Diyagramda "alt" fragment ile bu yetki kontrolü modellenmiştir. Bu, RBAC (Role-Based Access Control) yaklaşımının pratik uygulamasıdır. NFR-02 güvenlik isterimizi doğrudan karşılar.
""")

# ─── SLIDE 13: DAR Backend ───
add_notes(slides[12], """
Proje Adımı 2'de hazırladığımız Teknoloji Karar Analizi yani DAR raporu kapsamında, her teknoloji alternatifini 1'den 5'e kadar puanladık.

Backend Framework Karar Matrisi'ne bakalım:

Django REST Framework (Python): Toplam 19 puan
- Geliştirme hızı 5/5: Çok az boilerplate kod yazılır, dahili admin paneli, ORM otomatik migration...
- Güvenlik 5/5: JWT, RBAC, CSRF koruması dahili olarak gelir.
- Ekip yetkinliği 5/5: Ekibimiz Python konusunda en deneyimli.
- OOP uyumu 4/5: Python nesne yönelimli ama Java kadar katı tip sistemi yok.

Express.js (Node): Toplam 12 puan
- Hızlı ama güvenlik katmanları sıfırdan yazılmalı.

Spring Boot (Java): Toplam 13 puan
- OOP uyumu en yüksek (5/5) ama öğrenme eğrisi çok dik ve ekip deneyimi düşük.

Sonuç: Güvenlik altyapısı ve ekibin hızı göz önüne alınarak Django REST Framework seçildi.
""")

# ─── SLIDE 14: DAR DB & Frontend ───
add_notes(slides[13], """
Veritabanı Karar Matrisi:

PostgreSQL: Toplam 20 puan ile açık ara birinci.
- ACID uyumu 5/5: Sağlık verileri kesinlikle tutarlı olmalı.
- İlişki modelleme 5/5: Sınıf diyagramındaki 1'e çok ilişkiler ForeignKey ile birebir eşlenir.
- Ölçeklenebilirlik 5/5: Production ortamında yüksek yük altında çalışabilir.
- Medikal güvenlik 5/5: Şema izolasyonu ile multi-tenant veri güvenliği sağlanır.

MongoDB elenme nedeni: Şemasız yapısı sağlık verilerinde veri bütünlüğü riskleri taşır.
SQLite elenme nedeni: Production ortamında ölçeklenemez, tek dosya tabanlıdır.

Frontend Framework Karar Matrisi:

React.js: Toplam 20 puan ile birinci.
- En kritik avantajı: Dijital Odontogram yani diş şeması bileşeni için DOM manipülasyonunda Virtual DOM sayesinde performans kaybı olmuyor.
- Component-based mimarisi, sınıf diyagramındaki Boundary sınıflarıyla doğrudan eşleşiyor.

Alt kutuda: Backend dil olarak Python, ORM olarak Django ORM seçildi. Sınıf diyagramı nesnelerini tablolara birebir eşler.
""")

# ─── SLIDE 15: API Kontratları ───
add_notes(slides[14], """
Proje Adımı 3 kapsamında Opsiyon A'yı seçtik: İç ve Dış Arayüzlerin Gerçeklenmesi.

Swagger yani drf-spectacular kullanarak tüm API endpointlerini canlı olarak dokümante ettik.

4 temel API kontratımızı inceliyoruz:

A) POST /api/auth/token/ - JWT ile oturum açma:
Request olarak kullanıcı adı ve parola gönderilir. Başarılı olursa access ve refresh token çifti döner. Hatalı girişte 401, eksik parametrede 400 hata kodu döner.

C) POST /api/appointments/ - Çakışma kontrollü randevu:
Hasta ID, doktor ID, tarih ve saat gönderilir. Çakışma varsa 400, yetkisiz ise 403, başarılı ise 201 döner.

D) GET /api/dashboard/today/ - Günlük özet:
Tek endpoint'ten bugünkü randevular, toplam, tamamlanan ve hasta sayısı döner.

Sağ alt köşede katmanlar arası iletişim açıklanıyor: Model, ViewSet ve Serializer JSON paketleriyle haberleşir. get_serializer_class ile create ve update işlemleri farklı serializer'lara yönlendirilir.
""")

# ─── SLIDE 16: Dış Servis & Multi-Tenant ───
add_notes(slides[15], """
Sol tarafta Harici SMS Entegrasyonu:
- Randevu oluşturulduğunda harici bir SMS API'sine bildirim isteği atılır.
- Payload: api_key, phone_number ve message içerir.
- ÖNEMLİ: SMS API timeout olursa veya hata verirse uygulama ÇÖKMEZ. Hata log sisteme kaydedilir ve randevu asenkron olarak tamamlanır. Bu "fail gracefully" prensibidir.
- Frontend tarafında axios kütüphanesi ile JWT token otomatik olarak her isteğin Header'ına eklenir (interceptor pattern).

Sağ tarafta Multi-Tenant Mimarisi:
- django-tenants kütüphanesi ile "schema-per-tenant" yaklaşımı kullanıyoruz.
- HeaderTenantMiddleware, her gelen HTTP isteğinde domain veya X-Tenant header'ından hangi kliniğe ait olduğunu çözer.
- PostgreSQL şemasını o kliniğe yönlendirir.
- Sonuç: Her klinik kendi veritabanı şemasında çalışır. Kiracılar arası veri sızıntısı mimari düzeyde imkânsızdır.
- Bu yapıyı birim testlerde de doğruladık - UT-006'dan UT-011'e kadar HeaderTenantMiddleware testleri bu mekanizmayı test eder.
""")

# ─── SLIDE 17: Test Ortamı ve Kritik Sınıflar ───
add_notes(slides[16], """
Proje Adımı 4'te birim test çalışması gerçekleştirdik.

Test Araçlarımız:
- pytest: Python'un endüstri standardı test framework'ü. Sade assert sözdizimi ve az boilerplate.
- pytest-django: Her test için izole, otomatik rollback'li veritabanı sağlar. Testler birbirini etkilemez.
- factory-boy + Faker: Hasta, Doktor, Tedavi gibi test nesnelerini tek satırda Türkçe veriyle üretir.
- unittest.mock / monkeypatch: Veritabanı ve dış bağımlılıkları taklit ederek metodu gerçekten izole eder.

5 Kritik Sınıf seçtik:
1. AppointmentCreateSerializer - validate(): Randevu çakışması engelleme. Mock kullanmıyor çünkü doğrudan Django ORM ile test ediyor.
2. HeaderTenantMiddleware - __call__(): Kiracı çözümleme. Mock kullanıyor çünkü gerçek PostgreSQL şema işlemlerini taklit ediyoruz.
3. RegisterClinicView / CheckDomainView: Yeni klinik kayıt akışı. Mock kullanıyor.
4. PatientSerializer - create/update: Hasta ve anamnez kaydı. Mock yok.
5. TreatmentSerializer - validate(): Mükerrer tedavi engelleme. Mock yok.

Test Metodolojisi: Her metot için Pozitif ve Negatif senaryo. Negatif senaryoda kodun hatayı reddetmesi PASSED sayılır.
""")

# ─── SLIDE 18: Test Sonuç Matrisi ───
add_notes(slides[17], """
Test koşum komutu: pytest test_serializers.py test_middleware.py test_register.py → 54 passed.

Tabloda 8 temsili senaryoyu görüyorsunuz. Tamamını açıklayalım:

UT-001: Aynı hekime aynı tarih ve saatte ikinci randevu ekleniyor → ValidationError fırlatılıyor → PASSED. Bu en kritik testimiz.

UT-003: Aynı hekime ama FARKLI saatte randevu → Kabul ediliyor → PASSED. Çakışma algoritmasının doğru çalıştığını doğrular.

UT-006: X-Tenant header ile gelen istekte doğru kiracıya yönlendirme → Mock ile test → PASSED.

UT-011: Public tenant tanımsızken hiçbir eşleşme olmaması → 500 döndürmeli → PASSED.

UT-014: Zorunlu alan eksik gönderildiğinde 400 hatası → PASSED.

UT-022: Telefon olmadan hasta kaydı → is_valid = False, phone hatası → PASSED.

UT-025: Aynı gün, aynı diş, aynı tedavi türü mükerrer eklenmesi → ValidationError → PASSED.

Toplam: 25 temsili senaryo (13 Pozitif + 12 Negatif), TÜMÜ PASSED.
3 test dosyasının tam koşumu: 54 / 54 PASSED.
""")

# ─── SLIDE 19: Bakım ve Gelecek Önerileri ───
add_notes(slides[18], """
Dersin Bakım ayağına uygun olarak 3 kategoride gelecek planlarımızı sunduk:

Refactoring İhtiyaçları:
1. Servis Katmanı ayrımı: Şu an iş kuralları ViewSet ve Serializer içinde. Bunları bağımsız bir Service Layer'a taşımak bakım maliyetini düşürecek.
2. Test kapsamı genişletme: Birim testlerin yanına entegrasyon ve E2E testleri eklenecek.
3. Merkezi hata yönetimi: RFC 7807 Problem Details formatında standart hata yanıtları.

Yeni Özellik Önerileri:
- Gerçek SMS/WhatsApp entegrasyonu
- Takvimde sürükle-bırak
- Dijital röntgen dosyaları
- Raporlama modülü
- Mobil uygulama (React Native)
- Çok dilli destek

Versiyon Güncelleme Planı:
v1.1 (2025 Q3): SMS + takvim drag&drop
v1.2 (2025 Q4): Raporlama + PDF çıktı
v2.0 (2026 Q1): Mobil + çok dilli
v2.1 (2026 Q2): Röntgen + gelişmiş analitik
""")

# ─── SLIDE 20: Sonuç ───
add_notes(slides[19], """
Sonuç olarak, Yaşca Diş Kliniği Yönetim Sistemi'ni yazılım yaşam döngüsünün tüm aşamalarını kapsayacak şekilde geliştirdik:

Adım 1'de 6 fonksiyonel gereksinimi UML ile modelledik - 6 sekans ve 2 sınıf diyagramı.
Adım 2'de DAR raporu ile teknoloji yığınını karar matrisleriyle objektif olarak seçtik.
Adım 3'te iç ve dış arayüzleri Swagger ile dokümante edip multi-tenant altyapıyı kodladık.
Adım 4'te 5 kritik sınıfın birim testlerini hem pozitif hem negatif senaryolarla doğruladık.

Sayılarla Yaşca:
- 54/54 birim test PASSED
- 6 fonksiyonel gereksinim
- 6 NFR + 3 karar matrisi
- 8 domain varlığı
- %0 kiracılar arası veri sızıntısı

Temel çıktılarımız: Uçtan uca test edilmiş modüller, standartlara uygun UML dokümantasyonu, teknik karar verme yetkinliği, ölçeklenebilir multi-tenant SaaS mimarisi ve Swagger ile canlı API dokümantasyonu.

Sorularınız varsa yanıtlamaktan memnuniyet duyarız. Teşekkür ederiz.
""")

# ─── Save ───
prs.save(OUT_PATH)
print(f"Sunum basariyla olusturuldu: {OUT_PATH}")
print(f"Toplam slayt: {len(prs.slides)}")
