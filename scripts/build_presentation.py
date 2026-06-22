# -*- coding: utf-8 -*-
"""
Yaşca — Final Proje Raporu sunum üretici (v3).
İçerik %100 gerçek rapordan (Grup4_Final_Proje_Raporu.docx) alınmıştır.
UML bölümünde gerçek diyagram görselleri (C:\\Users\\CIHAN\\Desktop\\UML),
ayrı bölümde canlı uygulama ekran görüntüleri (C:\\Users\\CIHAN\\Desktop\\Yasca_Screenshots)
gömülür.

Çalıştır:  python scripts/build_presentation.py
Çıktı:     <Masaüstü>/Yasca_Final_Proje_Raporu.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------- renkler / stil
NAVY      = RGBColor(0x0E, 0x51, 0x89)
NAVY_DK   = RGBColor(0x0C, 0x42, 0x70)
CREAM     = RGBColor(0xF3, 0xEF, 0xE0)
CREAM_LT  = RGBColor(0xFA, 0xF7, 0xEC)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
INK       = RGBColor(0x2B, 0x2B, 0x2B)
MUTED     = RGBColor(0x5A, 0x5A, 0x5A)
GREEN     = RGBColor(0x1F, 0x9D, 0x55)
TEAL      = RGBColor(0x1F, 0xB6, 0x8D)
GOLD      = RGBColor(0xB8, 0x86, 0x12)
BORDER    = RGBColor(0xCC, 0xCC, 0xCC)
LIGHTBLUE = RGBColor(0xE8, 0xEF, 0xF6)
SERIF = "Georgia"
SANS  = "Calibri"

# ---------------------------------------------------------------- görsel yolları
UML = r"C:\Users\CIHAN\Desktop\UML"
IMG_CLASS_DOMAIN = os.path.join(UML, "WhatsApp Image 2026-04-25 at 17.27.36.jpeg")
IMG_CLASS_APP    = os.path.join(UML, "WhatsApp Image 2026-04-25 at 17.28.41.jpeg")
IMG_SD02         = os.path.join(UML, "WhatsApp Image 2026-04-25 at 17.24.32.jpeg")
IMG_SD03         = os.path.join(UML, "WhatsApp Image 2026-04-25 at 17.06.03.jpeg")
IMG_SD04         = os.path.join(UML, "WhatsApp Image 2026-04-25 at 17.23.21.jpeg")
IMG_SD06         = os.path.join(UML, "WhatsApp Image 2026-04-25 at 17.25.28.jpeg")
SHOT = r"C:\Users\CIHAN\Desktop\Yasca_Screenshots"
SH_LANDING  = os.path.join(SHOT, "01_landing.png")
SH_REGISTER = os.path.join(SHOT, "02_register.png")
SH_LOGIN    = os.path.join(SHOT, "03_login.png")
SH_DASH     = os.path.join(SHOT, "10_dashboard.png")
SH_PATIENTS = os.path.join(SHOT, "11_hastalar.png")
SH_APPTS    = os.path.join(SHOT, "12_randevular.png")
SH_PROFILE  = os.path.join(SHOT, "13_hasta_profil.png")

EMU_W = Inches(13.333); EMU_H = Inches(7.5)
prs = Presentation(); prs.slide_width = EMU_W; prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------- yardımcılar
def add_slide(bg=CREAM):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background(); r.shadow.inherit = False
    sp = r._element; sp.getparent().remove(sp); s.shapes._spTree.insert(2, sp)
    return s

def _set_font(run, size, color, bold=False, italic=False, font=SANS):
    run.font.size = Pt(size); run.font.color.rgb = color
    run.font.bold = bold; run.font.italic = italic; run.font.name = font

def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2); tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    return tb, tf

def para(tf, text, size, color, bold=False, italic=False, font=SANS,
         align=PP_ALIGN.LEFT, space_after=6, space_before=0, level=0, bullet=None, line=1.05, first=False):
    p = tf.paragraphs[0] if (first and tf.paragraphs[0].text == "" and not tf.paragraphs[0].runs) else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(space_after); p.space_before = Pt(space_before)
    p.line_spacing = line; p.level = level
    if bullet:
        r1 = p.add_run(); r1.text = bullet + "  "; _set_font(r1, size, color, bold, italic, font)
        r2 = p.add_run(); r2.text = text; _set_font(r2, size, color, bold, italic, font)
    else:
        r = p.add_run(); r.text = text; _set_font(r, size, color, bold, italic, font)
    return p

def section_title(slide, text, sub=None):
    tb, tf = textbox(slide, Inches(0.6), Inches(0.32), Inches(12.1), Inches(0.85))
    para(tf, text, 30, NAVY_DK, bold=True, font=SERIF, first=True, space_after=0)
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(1.12), Inches(12.1), Pt(3))
    ln.fill.solid(); ln.fill.fore_color.rgb = NAVY_DK; ln.line.fill.background(); ln.shadow.inherit = False
    if sub:
        tb2, tf2 = textbox(slide, Inches(0.62), Inches(1.17), Inches(12.1), Inches(0.36))
        para(tf2, sub, 14, MUTED, italic=True, first=True)

def card(slide, l, t, w, h, fill=CREAM_LT, line=RGBColor(0xDD, 0xD6, 0xC0), radius=0.05):
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    c.fill.solid(); c.fill.fore_color.rgb = fill
    c.line.color.rgb = line; c.line.width = Pt(1); c.shadow.inherit = False
    try: c.adjustments[0] = radius
    except Exception: pass
    return c

def card_content(slide, l, t, w, h, icon, title, body_lines, fill=CREAM_LT,
                 title_color=NAVY_DK, body_color=INK, icon_color=NAVY, title_size=18, body_size=13):
    card(slide, l, t, w, h, fill=fill)
    pad = Inches(0.26)
    tb, tf = textbox(slide, l + pad, t + Inches(0.2), w - 2 * pad, h - Inches(0.36))
    if icon:
        para(tf, icon, 22, icon_color, first=True, space_after=4)
        para(tf, title, title_size, title_color, bold=True, font=SERIF, space_after=8)
    else:
        para(tf, title, title_size, title_color, bold=True, font=SERIF, first=True, space_after=8)
    for ln in body_lines:
        if isinstance(ln, tuple):
            txt, b = ln
            para(tf, txt, body_size, body_color, bullet=b, space_after=4, line=1.05)
        else:
            para(tf, ln, body_size, body_color, space_after=4, line=1.05)
    return tf

def IX(v): return Inches(v)

def make_table(slide, rows, l, t, w, h, col_w, font_size=10, header_size=11,
               status_col=None, first_bold=True):
    nr = len(rows); nc = len(rows[0])
    tbl = slide.shapes.add_table(nr, nc, l, t, w, h).table
    for i, cw in enumerate(col_w):
        tbl.columns[i].width = Inches(cw)
    for ci in range(nc):
        tbl.cell(0, ci).fill.solid(); tbl.cell(0, ci).fill.fore_color.rgb = NAVY
    for ri in range(1, nr):
        band = CREAM_LT if ri % 2 else WHITE
        # son satır (Toplam) vurgusu opsiyonel değil; normal bandlama
        for ci in range(nc):
            tbl.cell(ri, ci).fill.solid(); tbl.cell(ri, ci).fill.fore_color.rgb = band
    for ri in range(nr):
        for ci in range(nc):
            cell = tbl.cell(ri, ci); cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Pt(5); cell.margin_right = Pt(4); cell.margin_top = Pt(1); cell.margin_bottom = Pt(1)
            p = cell.text_frame.paragraphs[0]; p.text = ""; p.line_spacing = 1.0
            r = p.add_run(); r.text = str(rows[ri][ci])
            if ri == 0:
                _set_font(r, header_size, WHITE, bold=True)
            elif status_col is not None and ci == status_col:
                _set_font(r, font_size, GREEN, bold=True)
            else:
                _set_font(r, font_size, INK, bold=(first_bold and ci == 0))
    return tbl

def image_card(slide, path, l, t, w, h, caption=None, border=True):
    cap_h = Inches(0.34) if caption else Inches(0.0)
    avail_h = h - cap_h
    if not path or not os.path.exists(path):
        card(slide, l, t, w, avail_h, fill=WHITE)
        tb, tf = textbox(slide, l, t, w, avail_h, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, "[görsel bulunamadı]", 12, MUTED, align=PP_ALIGN.CENTER, first=True)
    else:
        pic = slide.shapes.add_picture(path, l, t, width=w)
        if pic.height > avail_h:
            pic._element.getparent().remove(pic._element)
            pic = slide.shapes.add_picture(path, l, t, height=int(avail_h))
        pic.left = l + int((w - pic.width) / 2)
        pic.top = t + int((avail_h - pic.height) / 2)
        if border:
            pic.line.color.rgb = BORDER; pic.line.width = Pt(0.75)
    if caption:
        tb, tf = textbox(slide, l, t + avail_h, w, cap_h)
        para(tf, caption, 11, MUTED, align=PP_ALIGN.CENTER, first=True)


# ================================================================ 1) BAŞLIK
s = add_slide(NAVY)
tb, tf = textbox(s, IX(1.0), IX(0.7), IX(11.3), IX(1.5), anchor=MSO_ANCHOR.TOP)
para(tf, "İSTANBUL SAĞLIK VE TEKNOLOJİ ÜNİVERSİTESİ", 16, RGBColor(0xCF,0xDD,0xEC), bold=True, align=PP_ALIGN.CENTER, first=True, space_after=2)
para(tf, "Mühendislik ve Doğa Bilimleri Fakültesi · Yazılım Mühendisliği Bölümü", 12.5, RGBColor(0x9E,0xB8,0xD2), align=PP_ALIGN.CENTER, space_after=2)
para(tf, "YAZ402 — Yazılım Gerçekleme, Test ve Bakım", 13, RGBColor(0xCF,0xDD,0xEC), italic=True, align=PP_ALIGN.CENTER)
tb, tf = textbox(s, IX(1.0), IX(2.45), IX(11.3), IX(2.0), anchor=MSO_ANCHOR.MIDDLE)
para(tf, "Final Proje Raporu", 52, CREAM, bold=True, font=SERIF, align=PP_ALIGN.CENTER, first=True, space_after=8)
para(tf, "Yaşca Diş Kliniği Yönetim Sistemi", 24, WHITE, bold=True, align=PP_ALIGN.CENTER, space_after=2)
para(tf, "Çok-Kiracılı (Multi-Tenant) SaaS Mimarisi", 16, RGBColor(0xBF,0xD2,0xE6), italic=True, align=PP_ALIGN.CENTER)
tb, tf = textbox(s, IX(1.0), IX(4.95), IX(11.3), IX(2.0), anchor=MSO_ANCHOR.TOP)
para(tf, "Grup 4", 17, CREAM, bold=True, align=PP_ALIGN.CENTER, first=True, space_after=6)
para(tf, "Yaman Halloum   •   Ali Üre   •   Cihan Kurtbey   •   Berkay Aydın", 15, RGBColor(0xCF,0xDD,0xEC), align=PP_ALIGN.CENTER, space_after=14)
para(tf, "Django REST Framework (Python)  •  React.js (TypeScript)  •  PostgreSQL", 12, RGBColor(0x9E,0xB8,0xD2), align=PP_ALIGN.CENTER, space_after=4)
para(tf, "2024 – 2025 Bahar Dönemi", 12.5, RGBColor(0xBF,0xD2,0xE6), align=PP_ALIGN.CENTER)


# ================================================================ 2) İÇİNDEKİLER
s = add_slide()
section_title(s, "İçindekiler")
items = [
    ("1", "Giriş ve Proje Kapsamı"), ("2", "Yazılım Gereksinimleri (FR / NFR)"),
    ("3", "Detaylı UML Tasarımı (Sekans + Sınıf)"), ("4", "Teknoloji Karar Analizi (DAR)"),
    ("5", "İç ve Dış Arayüzler"), ("6", "Birim Test Sonuçları"),
    ("7", "Bakım ve Gelecek Önerileri"), ("8", "Sonuç + Canlı Uygulama"),
]
colw = IX(5.95); rowh = 1.15
for i, (no, title) in enumerate(items):
    col = i % 2; row = i // 2
    x = IX(0.6 + col * 6.25); y = IX(1.7 + row * 1.32)
    card(s, x, y, colw, IX(1.1), fill=WHITE)
    bub = s.shapes.add_shape(MSO_SHAPE.OVAL, x + IX(0.22), y + IX(0.3), IX(0.5), IX(0.5))
    bub.fill.solid(); bub.fill.fore_color.rgb = NAVY; bub.line.fill.background(); bub.shadow.inherit = False
    bp = bub.text_frame.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
    br = bp.add_run(); br.text = no; _set_font(br, 18, WHITE, bold=True, font=SERIF)
    tb, tf = textbox(s, x + IX(0.95), y, colw - IX(1.1), IX(1.1), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, title, 15, NAVY_DK, bold=True, font=SERIF, first=True)


# ================================================================ 3) GİRİŞ VE KAPSAM
s = add_slide()
section_title(s, "1. Giriş ve Proje Kapsamı")
tb, tf = textbox(s, IX(0.6), IX(1.5), IX(6.05), IX(5.6))
para(tf, "Projenin Amacı", 18, NAVY_DK, bold=True, font=SERIF, first=True, space_after=6)
para(tf, "Diş kliniklerinin günlük operasyonlarını dijitalleştiren ve merkezi bir platformdan "
         "yöneten çok-kiracılı (multi-tenant) bir SaaS uygulaması. Hasta kaydı, randevu, tedavi/finans, "
         "dijital odontogram ve klinik yönetimini tek çatıda sunar. Her kliniğin verisi PostgreSQL "
         "şema tabanlı mimari ile fiziksel olarak izole edilir.", 12.5, INK, space_after=12, line=1.13)
para(tf, "Ana Hedefler", 14, NAVY_DK, bold=True, space_after=4)
for t in ["Hasta kayıt & anamnez yönetimini dijitalleştirmek",
          "Çakışma kontrollü randevu ile verimliliği artırmak",
          "Tedavi + ödeme takibini entegre sunmak",
          "Dijital diş şeması (Odontogram) ile görselleştirme",
          "JWT kimlik doğrulama + rol bazlı yetkilendirme (RBAC)"]:
    para(tf, t, 11.5, MUTED, bullet="•", space_after=2, line=1.05)

card_content(s, IX(6.9), IX(1.5), IX(5.85), IX(2.55), None, "Hedef Kullanıcılar (Roller)",
             [("Yönetici (Admin): klinik ayarları, çalışma saatleri, personel, tüm veri", "•"),
              ("Hekim (Doktor): kendi hastaları, tedavi girişi, odontogram, dashboard", "•"),
              ("Asistan: hasta kaydı, randevu oluşturma, ödeme girişi", "•")],
             fill=WHITE, body_size=12, title_size=16)
card_content(s, IX(6.9), IX(4.2), IX(5.85), IX(2.9), None, "Kritik Başarı Kriterleri",
             [("Randevu çakışması %0 olmalı", "✓"),
              ("Kiracılar arası veri sızıntısı kesinlikle önlenmeli", "✓"),
              ("API yanıt süresi < 2 saniye", "✓"),
              ("En az 54 birim testin tamamı geçmeli", "✓"),
              ("Tüm API'ler Swagger (OpenAPI) ile dokümante", "✓")],
             fill=NAVY, title_color=CREAM, body_color=LIGHTBLUE, body_size=12, title_size=16)


# ================================================================ 4) FONKSİYONEL İSTERLER
s = add_slide()
section_title(s, "2.1. Fonksiyonel İsterler (FR)")
rows = [
    ("ID", "Gereksinim", "Açıklama"),
    ("FR-01", "Kimlik Doğrulama", "Klinik personelinin JWT tabanlı güvenli girişi; access + refresh token çifti döndürülmesi."),
    ("FR-02", "Hasta Yönetimi", "Yeni hasta kaydı ve tıbbi özgeçmiş (Anamnez) eklenmesi; bilgilerin güncellenmesi ve listelenmesi."),
    ("FR-03", "Randevu Yönetimi", "Çakışma kontrollü randevu oluşturma; aynı hekime aynı tarih/saatte ikinci randevu engellenir."),
    ("FR-04", "Klinik Operasyon (Dashboard)", "Günlük özet: bugünün randevuları, tamamlanan işlem sayısı ve aktif hasta sayısı."),
    ("FR-05", "Tedavi ve Finans", "Tedavilerin girilmesi, ödeme kayıtları ve tedavi-finans entegrasyonu."),
    ("FR-06", "Sistem Yönetimi", "Yöneticinin klinik çalışma saatleri/ayarlarını güncellemesi; 403 ile yalnızca Admin erişimi."),
]
make_table(s, rows, IX(0.6), IX(1.55), IX(12.13), IX(4.9), [0.95, 2.85, 8.33], font_size=11.5, header_size=12)


# ================================================================ 5) FONKSİYONEL OLMAYAN İSTERLER
s = add_slide()
section_title(s, "2.2. Fonksiyonel Olmayan İsterler (NFR)")
rows = [
    ("ID", "Kategori", "Açıklama"),
    ("NFR-01", "Performans", "API sorgu süresi 2 saniyeden az olmalıdır."),
    ("NFR-02", "Güvenlik", "JWT tabanlı kimlik doğrulama; RBAC (rol bazlı erişim) ile yetkilendirme."),
    ("NFR-03", "Veri İzolasyonu", "Her klinik verisi PostgreSQL şema seviyesinde fiziksel olarak izole edilmelidir."),
    ("NFR-04", "Ölçeklenebilirlik", "Birden fazla kliniğin eşzamanlı kullanımı desteklenmelidir."),
    ("NFR-05", "Bakım Yapılabilirlik", "Modüler mimari ve yüksek test kapsamı ile bakım maliyeti düşürülmelidir."),
    ("NFR-06", "Kullanılabilirlik", "Kullanıcı arayüzü sezgisel ve Türkçe olmalıdır."),
]
make_table(s, rows, IX(0.6), IX(1.55), IX(12.13), IX(4.7), [1.15, 2.85, 8.13], font_size=12, header_size=12)


# ================================================================ 6) UML GENEL BAKIŞ
s = add_slide()
section_title(s, "3. Detaylı UML Tasarımı — Genel Bakış")
tb, tf = textbox(s, IX(0.6), IX(1.45), IX(12.1), IX(0.7))
para(tf, "Sistem, Boundary (Arayüz) – Control (Servis/API) – Entity (Varlık) ayrımına göre standart UML "
         "notasyonuyla modellendi. 6 sekans diyagramı (SD-01..06) ve 2 sınıf diyagramı (Domain + Uygulama) tasarlandı.",
     12.5, INK, first=True, line=1.12)
card_content(s, IX(0.6), IX(2.35), IX(5.95), IX(4.5), None, "Sekans Diyagramları (Dinamik)",
             [("SD-01 — Kullanıcı Girişi ve Oturum Doğrulaması (FR-01)", "•"),
              ("SD-02 — Hasta Kaydı ve Anamnez Güncelleme (FR-02)", "•"),
              ("SD-03 — Randevu Oluşturma + Çakışma Kontrolü (FR-03)", "•"),
              ("SD-04 — Günlük Dashboard Özet Akışı (FR-04)", "•"),
              ("SD-05 — Tedavi ve Ödeme Girişi (FR-05)", "•"),
              ("SD-06 — Klinik Ayarları Güncelleme (FR-06)", "•")],
             fill=WHITE, body_size=12.5, title_size=16)
card_content(s, IX(6.8), IX(2.35), IX(5.95), IX(4.5), None, "Sınıf Diyagramları (Statik)",
             [("Domain Katmanı: 8 varlık (Clinic, CustomUser, Patient,", "•"),
              ("   Anamnesis, Appointment, Treatment, TreatmentType, Payment)", None),
              ("Uygulama Katmanı: Boundary (UI) + Control (ViewSet/Serializer)", "•"),
              ("İlişkiler: 1..* Association, 1..1 / 1..* Composition, Aggregation", "•"),
              ("Tutarlılık: sekans mesajları = ApiService metodları (birebir)", "•")],
             fill=NAVY, title_color=CREAM, body_color=LIGHTBLUE, body_size=12.5, title_size=16)


# ================================================================ 7) SINIF — DOMAIN
s = add_slide()
section_title(s, "3.2.1. Sınıf Diyagramı — Domain Katmanı", sub="Veri varlıkları ve ilişkileri (1..* Association / Composition / Aggregation)")
image_card(s, IM_DOMAIN := IMG_CLASS_DOMAIN, IX(0.55), IX(1.65), IX(7.4), IX(5.5))
tb, tf = textbox(s, IX(8.15), IX(1.7), IX(4.65), IX(5.4))
para(tf, "Sınıflar Arası İlişkiler", 15, NAVY_DK, bold=True, font=SERIF, first=True, space_after=7)
for t in ["Clinic → Patient: 1..* Association",
          "Clinic → CustomUser: 1..* Association",
          "Patient → Appointment: 1..* Composition",
          "Patient → Anamnesis: 1..1 Composition",
          "Appointment → Treatment: 1..* Aggregation",
          "Treatment → Payment: 1..1 Association"]:
    para(tf, t, 12, INK, bullet="•", space_after=6, line=1.1)
para(tf, "Temel varlıklar: Clinic, CustomUser, Patient, Anamnesis, Appointment, "
         "Treatment, TreatmentType, Payment.", 11, MUTED, italic=True, space_before=4, line=1.12)


# ================================================================ 8) SINIF — UYGULAMA
s = add_slide()
section_title(s, "3.2.2. Sınıf Diyagramı — Uygulama Katmanı", sub="Boundary (Frontend) ↔ ApiService (Control) ↔ Backend ViewSet / Serializer")
image_card(s, IMG_CLASS_APP, IX(0.55), IX(1.65), IX(8.2), IX(5.5))
tb, tf = textbox(s, IX(8.95), IX(1.7), IX(3.85), IX(5.4))
para(tf, "Stereotipler", 15, NAVY_DK, bold=True, font=SERIF, first=True, space_after=6)
para(tf, "Boundary (UI)", 12.5, NAVY_DK, bold=True, space_after=2)
para(tf, "LoginPage, Dashboard, PatientDialog, PatientProfile, AppointmentDialog, ClinicSettingsPage", 11, MUTED, space_after=8, line=1.12)
para(tf, "Control (API)", 12.5, NAVY_DK, bold=True, space_after=2)
para(tf, "ApiService, Auth/CurrentUserView, PatientViewSet, AppointmentViewSet (+Serializer), DashboardView", 11, MUTED, space_after=8, line=1.12)
para(tf, "Sorumlulukların ayrımı (SoC): Boundary, Entity ile doğrudan konuşmaz; araya Control katmanı girer (MVC/MVT).", 11, INK, italic=True, line=1.12)


# ================================================================ 9-12) SEKANS DİYAGRAMLARI (gerçek görseller)
def seq_slide(title, sub, img):
    s = add_slide()
    section_title(s, title, sub=sub)
    image_card(s, img, IX(0.7), IX(1.7), IX(11.93), IX(5.4))
    return s

seq_slide("3.1. Sekans Diyagramı — SD-02", "Hasta Kaydı ve Anamnez Güncelleme (FR-02): PatientDialog → ApiService → PatientViewSet → Patient + Anamnesis", IMG_SD02)
seq_slide("3.1. Sekans Diyagramı — SD-03", "Randevu Oluşturma + Çakışma Kontrolü (FR-03): exists(doctor,date,time) → alt [çakışma var → 400 / yok → 201]", IMG_SD03)
seq_slide("3.1. Sekans Diyagramı — SD-04", "Günlük Dashboard Özet Akışı (FR-04): bugünün randevuları + tamamlanan işlem + aktif hasta sayısı", IMG_SD04)
seq_slide("3.1. Sekans Diyagramı — SD-06", "Klinik Ayarları Güncelleme (FR-06): yetki kontrolü → alt [admin değil → 403 Forbidden / admin → 200 OK]", IMG_SD06)


# ================================================================ 13) DAR — Backend
s = add_slide()
section_title(s, "4. Teknoloji Karar Analizi (DAR) — Backend")
tb, tf = textbox(s, IX(0.6), IX(1.45), IX(12.1), IX(1.05))
para(tf, "Dil: Python — nesne yönelimli yapı + hızlı prototipleme + zengin güvenlik/şifreleme kütüphaneleri. "
         "Framework: Django REST — MVC mimarisi sekans diyagramlarındaki kontrolcülerle örtüşür; yerleşik JWT "
         "katmanı (Express'te sıfırdan yazılması gerekirdi).", 12.5, INK, first=True, line=1.13)
rows = [
    ("Kriter", "Django REST (Python)", "Express.js (Node)", "Spring Boot (Java)"),
    ("Geliştirme Hızı / Öğrenme Eğrisi", "5", "4", "2"),
    ("Dahili Güvenlik (Auth/RBAC)", "5", "2", "4"),
    ("Ekip Yetkinliği", "5", "3", "2"),
    ("Sınıf Diyagramı (OOP) Uyumu", "4", "3", "5"),
    ("Toplam Puan", "19", "12", "13"),
]
make_table(s, rows, IX(1.3), IX(2.75), IX(10.7), IX(3.4), [4.4, 2.2, 2.05, 2.05], font_size=12.5, header_size=12)
tb, tf = textbox(s, IX(0.6), IX(6.5), IX(12.1), IX(0.5))
para(tf, "Sonuç: Güvenlik altyapısı ve ekip hızı nedeniyle Django REST Framework seçildi (19/20).", 12.5, NAVY_DK, bold=True, first=True)


# ================================================================ 14) DAR — Veritabanı + Frontend
s = add_slide()
section_title(s, "4. DAR — Veritabanı ve Frontend Karar Matrisleri")
tb, tf = textbox(s, IX(0.6), IX(1.35), IX(6.0), IX(0.35))
para(tf, "Veri Tabanı", 14, NAVY_DK, bold=True, font=SERIF, first=True)
rows = [
    ("Kriter", "PostgreSQL", "MongoDB", "SQLite"),
    ("Veri Bütünlüğü (ACID)", "5", "2", "4"),
    ("Sınıf/İlişki Modelleme", "5", "3", "5"),
    ("Üretim Ölçeklenebilirliği", "5", "5", "1"),
    ("Medikal Veri Güvenliği", "5", "3", "2"),
    ("Toplam Puan", "20", "13", "12"),
]
make_table(s, rows, IX(0.6), IX(1.7), IX(6.0), IX(2.7), [2.7, 1.2, 1.1, 1.0], font_size=10.5, header_size=10.5)
tb, tf = textbox(s, IX(0.6), IX(4.5), IX(6.0), IX(0.5))
para(tf, "Sonuç: İlişkisel yapı + ACID zorunluluğu → PostgreSQL (20).", 11.5, NAVY_DK, bold=True, first=True, line=1.1)

tb, tf = textbox(s, IX(6.95), IX(1.35), IX(6.0), IX(0.35))
para(tf, "Frontend", 14, NAVY_DK, bold=True, font=SERIF, first=True)
rows = [
    ("Kriter", "React.js", "Vue.js", "Angular"),
    ("Topluluk Desteği / Kaynak", "5", "4", "3"),
    ("Odontogram DOM Performansı", "5", "4", "3"),
    ("Ekip Yetkinliği", "5", "2", "1"),
    ("Component (Boundary) Uyumu", "5", "5", "5"),
    ("Toplam Puan", "20", "15", "12"),
]
make_table(s, rows, IX(6.95), IX(1.7), IX(6.0), IX(2.7), [2.9, 1.05, 1.05, 1.0], font_size=10.5, header_size=10.5)
tb, tf = textbox(s, IX(6.95), IX(4.5), IX(6.0), IX(0.5))
para(tf, "Sonuç: DOM gücü + ekip tecrübesi → React.js (20).", 11.5, NAVY_DK, bold=True, first=True, line=1.1)

tb, tf = textbox(s, IX(0.6), IX(5.4), IX(12.2), IX(1.4))
para(tf, "Seçilen Yığın:  Django REST Framework (Python)  ·  React.js (TypeScript)  ·  PostgreSQL", 14, NAVY_DK, bold=True, font=SERIF, first=True, space_after=4)
para(tf, "ORM: Django ORM — sınıf diyagramındaki nesneler tablolara eşlenir; \"Bire-Çok\" ilişkiler ForeignKey ile koda dökülür.", 12, MUTED, italic=True, line=1.12)


# ================================================================ 15) ARAYÜZLER — API KONTRATLARI
s = add_slide()
section_title(s, "5. İç ve Dış Arayüzler — API Kontratları", sub="OpenAPI / Swagger (drf-spectacular) ile canlı dokümante edilmiş RESTful uç noktalar")
card_content(s, IX(0.6), IX(1.7), IX(5.95), IX(2.4), None, "POST /api/auth/token/  (FR-01)",
             [("Request: {username, password}", "•"),
              ("Response 200: {access, refresh} (JWT)", "•"),
              ("Hata: 401 Unauthorized · 400 Bad Request", "•")], fill=WHITE, body_size=12, title_size=15)
card_content(s, IX(6.8), IX(1.7), IX(5.95), IX(2.4), None, "POST /api/appointments/  (FR-03)",
             [("Request: {patient, doctor, date, time}", "•"),
              ("Response 201: oluşturulan randevu (id)", "•"),
              ("Hata: 400 (Çakışma) · 403 Forbidden", "•")], fill=WHITE, body_size=12, title_size=15)
card_content(s, IX(0.6), IX(4.25), IX(5.95), IX(2.65), None, "GET /api/dashboard/today/  (FR-04)",
             [("today_appointments: [...]", "•"),
              ("today_total: 12 · today_completed: 4", "•"),
              ("total_patients: 150", "•"),
              ("Tek uç noktadan günlük özet (tek istek)", "•")], fill=WHITE, body_size=12, title_size=15)
card_content(s, IX(6.8), IX(4.25), IX(5.95), IX(2.65), None, "GET & POST /api/patients/  (FR-02)",
             [("POST: {first_name, last_name, phone}", "•"),
              ("GET 200: [{id, full_name, phone, tckn, last_visit}]", "•"),
              ("Hata: 400 (zorunlu alan eksik)", "•"),
              ("Katmanlar JSON paketleriyle haberleşir", "•")], fill=NAVY, title_color=CREAM, body_color=LIGHTBLUE, body_size=12, title_size=15)


# ================================================================ 16) MULTI-TENANT + DIŞ SERVİS
s = add_slide()
section_title(s, "5.3. Multi-Tenant Mimarisi ve Dış Servisler")
card_content(s, IX(0.6), IX(1.7), IX(5.95), IX(5.1), None, "Multi-Tenant (django-tenants)",
             [("HeaderTenantMiddleware: HTTP isteğindeki domain/tenant", "•"),
              ("   kimliğini okur ve PostgreSQL şemasını o kliniğe yönlendirir", None),
              ("Veri güvenliği fiziksel izolasyonla sağlanır", "•"),
              ("Kiracılar arası veri sızıntısı mimari düzeyde engellenir", "•"),
              ("Frontend ↔ Backend: axios; JWT (Bearer) otomatik header'a eklenir", "•")],
             fill=WHITE, body_size=12.5, title_size=16)
card_content(s, IX(6.8), IX(1.7), IX(5.95), IX(5.1), None, "Dış Servis — SMS Entegrasyonu",
             [("Randevu oluşturulunca harici SMS API'sine istek atılır", "•"),
              ("Payload: {api_key, phone_number, message}", "•"),
              ("Hata yönetimi: dış API yanıt vermezse (Timeout)", "•"),
              ("   işlem çökmez; hata loglanır, randevu asenkron tamamlanır", None),
              ("Frontend baseURL: http://127.0.0.1:8000/api/ (geliştirme)", "•")],
             fill=NAVY, title_color=CREAM, body_color=LIGHTBLUE, body_size=12.5, title_size=16)


# ================================================================ 17) TEST — ORTAM + KRİTİK SINIFLAR
s = add_slide()
section_title(s, "6. Birim Test — Ortam ve Kritik Sınıflar")
tb, tf = textbox(s, IX(0.6), IX(1.35), IX(6.0), IX(0.35))
para(tf, "Test Ortamı ve Araçlar", 14, NAVY_DK, bold=True, font=SERIF, first=True)
rows = [
    ("Araç", "Görevi"),
    ("pytest", "Test koşum motoru (endüstri standardı)"),
    ("pytest-django", "İzole, otomatik rollback'li DB (django_db)"),
    ("factory-boy + Faker", "Türkçe sahte test verisi üretimi"),
    ("unittest.mock / monkeypatch", "Bağımlılık taklidi (gerçek izolasyon)"),
]
make_table(s, rows, IX(0.6), IX(1.7), IX(6.05), IX(2.5), [2.55, 3.5], font_size=10.5, header_size=11)
tb, tf = textbox(s, IX(6.95), IX(1.35), IX(6.0), IX(0.35))
para(tf, "Test Edilen 5 Kritik Sınıf", 14, NAVY_DK, bold=True, font=SERIF, first=True)
rows = [
    ("#", "Sınıf / Metot", "Mock"),
    ("1", "AppointmentCreateSerializer.validate()", "Hayır"),
    ("2", "HeaderTenantMiddleware.__call__()", "Evet"),
    ("3", "RegisterClinicView / CheckDomainView", "Evet"),
    ("4", "PatientSerializer.create/update()", "Hayır"),
    ("5", "TreatmentSerializer.validate()", "Hayır"),
]
make_table(s, rows, IX(6.95), IX(1.7), IX(6.0), IX(2.95), [0.5, 4.4, 1.1], font_size=10.5, header_size=11)
tb, tf = textbox(s, IX(0.6), IX(5.05), IX(12.2), IX(1.9))
para(tf, "Test Tasarım Metodolojisi", 14, NAVY_DK, bold=True, font=SERIF, first=True, space_after=4)
para(tf, "Pozitif senaryo: doğru girdide doğru sonuç (ör. boş randevu saatinin kabulü).", 12, INK, bullet="•", space_after=3, line=1.1)
para(tf, "Negatif senaryo: hatalı girdide sistemin hatayı fırlatması/işlemi reddetmesi (ör. çakışan randevu → ValidationError).", 12, INK, bullet="•", space_after=3, line=1.1)
para(tf, "Önemli: Negatif senaryoda kodun doğru reddetmesi başarı sayılır → durum PASSED.", 12, GOLD, bold=True, italic=True, line=1.1)


# ================================================================ 18) TEST SONUÇ MATRİSİ
s = add_slide()
section_title(s, "6.3. Birim Test Sonuç Matrisi")
tb, tf = textbox(s, IX(0.6), IX(1.3), IX(12.1), IX(0.5))
para(tf, "Koşum: pytest test_serializers.py test_middleware.py test_register.py → 54 passed. Temsilî satırlar:",
     12, INK, first=True, line=1.05)
rows = [
    ("Test ID", "Hedef Sınıf", "Metot", "Tür", "Senaryo", "Durum"),
    ("UT-001", "AppointmentCreateSerializer", "validate", "Negatif", "Aynı hekim/saat ikinci randevu reddi", "PASSED"),
    ("UT-003", "AppointmentCreateSerializer", "validate", "Pozitif", "Farklı saatte randevu kabul", "PASSED"),
    ("UT-006", "HeaderTenantMiddleware", "__call__", "Pozitif", "X-Tenant ile doğru kiracı çözümü", "PASSED"),
    ("UT-011", "HeaderTenantMiddleware", "__call__", "Negatif", "Public tenant yoksa HTTP 500", "PASSED"),
    ("UT-012", "RegisterClinicView", "post", "Pozitif", "Geçerli veriyle yeni klinik kaydı", "PASSED"),
    ("UT-014", "RegisterClinicView", "post", "Negatif", "Zorunlu alan eksik → HTTP 400", "PASSED"),
    ("UT-020", "PatientSerializer", "create", "Pozitif", "Nested anamnez ile hasta oluşturma", "PASSED"),
    ("UT-022", "PatientSerializer", "create", "Negatif", "Telefon olmadan kayıt reddi", "PASSED"),
    ("UT-024", "TreatmentSerializer", "validate", "Pozitif", "Aynı gün farklı diş → kabul", "PASSED"),
    ("UT-025", "TreatmentSerializer", "validate", "Negatif", "Aynı gün/diş mükerrer tedavi reddi", "PASSED"),
]
make_table(s, rows, IX(0.6), IX(1.95), IX(12.13), IX(4.1), [1.25, 3.35, 1.25, 1.2, 4.0, 1.08],
           font_size=9.5, header_size=10.5, status_col=5)
tb, tf = textbox(s, IX(0.6), IX(6.25), IX(12.1), IX(0.9))
para(tf, "Özet: 25 temsilî senaryo (13 Pozitif + 12 Negatif) — TÜMÜ PASSED · İlgili 3 test dosyasının tam koşumu: 54/54 passed.",
     13, NAVY_DK, bold=True, first=True)


# ================================================================ 19) BAKIM VE GELECEK
s = add_slide()
section_title(s, "7. Bakım ve Gelecek Önerileri")
card_content(s, IX(0.6), IX(1.65), IX(3.95), IX(3.3), "✎", "Refactoring İhtiyaçları",
             [("İş mantığını bağımsız Service Layer'a taşımak", "•"),
              ("Entegrasyon + E2E testleriyle kapsamı genişletmek", "•"),
              ("Merkezi hata yönetimi (RFC 7807 Problem Details)", "•")],
             fill=WHITE, body_size=11.5, title_size=15)
card_content(s, IX(4.69), IX(1.65), IX(3.95), IX(3.3), "✚", "Yeni Özellik Önerileri",
             [("SMS/WhatsApp ile otomatik randevu hatırlatma", "•"),
              ("Takvimde sürükle-bırak (drag & drop)", "•"),
              ("Dijital röntgen/görüntü dosyaları", "•"),
              ("Raporlama: gelir, hasta, tedavi istatistikleri", "•"),
              ("Mobil uygulama (React Native) · i18n", "•")],
             fill=WHITE, body_size=11.5, title_size=15)
card(s, IX(8.78), IX(1.65), IX(3.95), IX(3.3), fill=NAVY)
tb, tf = textbox(s, IX(9.0), IX(1.85), IX(3.55), IX(3.0))
para(tf, "⟳  Versiyon Planı", 15, CREAM, bold=True, font=SERIF, first=True, space_after=8)
for v, feat in [("v1.1", "SMS, takvim drag&drop — 2025 Q3"),
                ("v1.2", "Raporlama, PDF çıktı — 2025 Q4"),
                ("v2.0", "Mobil uygulama, çok dil — 2026 Q1"),
                ("v2.1", "Röntgen, gelişmiş analitik — 2026 Q2")]:
    para(tf, v + " — " + feat, 11.5, LIGHTBLUE, bullet="•", space_after=6, line=1.1)
# alt şerit
card(s, IX(0.6), IX(5.15), IX(12.13), IX(1.75), fill=CREAM_LT)
tb, tf = textbox(s, IX(0.85), IX(5.3), IX(11.7), IX(1.5))
para(tf, "Bakım odağı", 13, NAVY_DK, bold=True, first=True, space_after=4)
para(tf, "Modüler mimari, yüksek test kapsamı ve standart hata yönetimi ile uzun vadeli bakım maliyetinin "
         "düşürülmesi hedeflenir; teknik borç düzenli olarak gözden geçirilir.", 12, INK, line=1.15)


# ================================================================ 20) CANLI — PUBLIC
s = add_slide()
section_title(s, "Canlı Uygulama — Public Sayfalar", sub="yasca-dental-clinic.vercel.app — landing, klinik kayıt ve tenant giriş")
image_card(s, SH_LANDING, IX(0.6), IX(1.7), IX(4.25), IX(5.3), caption="Landing (tanıtım + fiyatlandırma)")
image_card(s, SH_REGISTER, IX(5.05), IX(1.7), IX(7.7), IX(2.55), caption="Klinik Kayıt — Hesabınızı Oluşturun")
image_card(s, SH_LOGIN, IX(5.05), IX(4.45), IX(7.7), IX(2.55), caption="Tenant Giriş — Kliniğinize Giriş Yapın")


# ================================================================ 21) CANLI — KLİNİK PANELİ 1
s = add_slide()
section_title(s, "Canlı Uygulama — Klinik Paneli (1)", sub="Giriş yapılmış görünüm: Dashboard (günlük özet) ve Hasta Yönetimi")
image_card(s, SH_DASH, IX(0.6), IX(1.7), IX(6.0), IX(5.2), caption="Dashboard — bugünkü randevular, bekleyen/toplam hasta")
image_card(s, SH_PATIENTS, IX(6.75), IX(1.7), IX(6.0), IX(5.2), caption="Hasta Yönetimi — arama + liste")


# ================================================================ 22) CANLI — KLİNİK PANELİ 2
s = add_slide()
section_title(s, "Canlı Uygulama — Klinik Paneli (2)", sub="Randevu Takvimi (haftalık) ve çok sekmeli Hasta Profili (Diş Şeması dahil)")
image_card(s, SH_APPTS, IX(0.6), IX(1.7), IX(6.0), IX(5.2), caption="Randevu Takvimi — haftalık görünüm")
image_card(s, SH_PROFILE, IX(6.75), IX(1.7), IX(6.0), IX(5.2), caption="Hasta Profili — Profil/Anamnez/Tedavi/Ödeme/Doküman/Diş Şeması")


# ================================================================ 23) SONUÇ
s = add_slide()
section_title(s, "8. Sonuç")
tb, tf = textbox(s, IX(0.6), IX(1.6), IX(6.5), IX(5.2))
para(tf, "\"Yaşca Diş Kliniği Yönetim Sistemi\" çok-kiracılı SaaS uygulaması; analiz, tasarım, gerçekleme, "
         "test ve bakım aşamalarının tamamını kapsayacak şekilde geliştirildi.", 13.5, INK, first=True, space_after=10, line=1.15)
para(tf, "Adım 1: 6 fonksiyonel gereksinim (FR-01..06) sekans + sınıf diyagramlarıyla (BCE) modellendi.", 12, MUTED, bullet="•", space_after=4, line=1.12)
para(tf, "Adım 2: DAR ile DRF + React + PostgreSQL yığını karar matrisleriyle seçildi.", 12, MUTED, bullet="•", space_after=4, line=1.12)
para(tf, "Adım 3: REST API'ler Swagger ile dokümante; frontend-backend + multi-tenant kodlandı.", 12, MUTED, bullet="•", space_after=4, line=1.12)
para(tf, "Adım 4: 5 kritik sınıf pozitif/negatif senaryolarla test edildi → 54/54 PASSED.", 12, MUTED, bullet="•", line=1.12)
card(s, IX(7.35), IX(1.65), IX(5.4), IX(5.15), fill=WHITE)
tb, tf = textbox(s, IX(7.65), IX(1.9), IX(4.85), IX(4.7))
para(tf, "Temel Çıktılar", 18, NAVY_DK, bold=True, font=SERIF, first=True, space_after=10)
for t in ["Uçtan uca test edilmiş modüller",
          "Standartlara uygun UML dokümantasyonu",
          "Teknik karar verme yetkinliği (DAR)",
          "Multi-tenant SaaS ile ölçeklenebilir altyapı",
          "Swagger ile canlı API dokümantasyonu",
          "54/54 birim testin başarıyla geçmesi"]:
    para(tf, t, 13, INK, bullet="✓", space_after=9, line=1.12)


# ---------------------------------------------------------------- konuşmacı notları
NOTES = [
    # 1 Başlık
    "Merhaba, herkese hoş geldiniz. Biz Grup 4 olarak — Yaman, Ali, Cihan ve Berkay — "
    "Yazılım Gerçekleme, Test ve Bakım dersi kapsamında geliştirdiğimiz 'Yaşca Diş Kliniği "
    "Yönetim Sistemi' projesini sunacağız. Yaşca, diş kliniklerinin tüm operasyonlarını tek "
    "platformdan yöneten, çok-kiracılı bir SaaS uygulaması. Sunumumuzda sistem analizinden "
    "başlayıp tasarım, teknoloji kararları, gerçekleme ve test sonuçlarına kadar tüm süreci ele alacağız.",
    # 2 İçindekiler
    "Sunum akışımız şu şekilde: önce projenin amacını ve gereksinimleri anlatacağız. Ardından UML "
    "tasarımımıza, yani sınıf ve sekans diyagramlarımıza bakacağız. Sonra hangi teknolojileri neden "
    "seçtiğimizi karar matrisleriyle göstereceğiz. Arayüzler ve API kontratlarının ardından birim test "
    "sonuçlarımızı, bakım planımızı ve son olarak da canlı uygulamadan ekran görüntülerini paylaşacağız.",
    # 3 Giriş
    "Sağlık sektöründe birçok klinik hâlâ kağıt ya da birbirinden kopuk yazılımlarla çalışıyor; bu da "
    "veri kaybı, randevu çakışması ve güvenlik açıklarına yol açıyor. Biz bu sorunları çözmek için "
    "Yaşca'yı geliştirdik. Her kliniğin verisi PostgreSQL şema seviyesinde fiziksel olarak izole ediliyor. "
    "Üç kullanıcı rolümüz var: Yönetici, Hekim ve Asistan. Başarı kriterlerimiz net: randevu çakışması "
    "sıfır, kiracılar arası sızıntı yok, yanıt süresi 2 saniyenin altında ve en az 54 birim test başarıyla geçmeli.",
    # 4 FR
    "Altı temel fonksiyonel gereksinimimiz var. FR-01 JWT ile güvenli giriş; FR-02 hasta ve anamnez "
    "yönetimi; FR-03 — projemizin kalbi — çakışma kontrollü randevu, yani aynı hekime aynı saatte ikinci "
    "randevu engelleniyor; FR-04 günlük dashboard özeti; FR-05 tedavi ve finans entegrasyonu; FR-06 ise "
    "yalnızca yöneticinin erişebildiği klinik ayarları. Her gereksinim, birazdan göreceğiniz diyagramlara doğrudan karşılık geliyor.",
    # 5 NFR
    "Bir de kalite, yani fonksiyonel olmayan gereksinimlerimiz var. Performansta hedef 2 saniyenin altı; "
    "güvenlikte JWT ve rol bazlı erişim; en kritiği veri izolasyonu — her klinik kendi şemasında. Ayrıca "
    "çok kiracılı yapıyla ölçeklenebilirlik, modüler mimariyle bakım kolaylığı ve sezgisel Türkçe arayüz.",
    # 6 UML overview
    "Tasarımı standart UML ile, Boundary–Control–Entity ayrımına göre yaptık; yani arayüz, servis ve veri "
    "katmanlarını net biçimde ayırdık. Altı kullanım durumu için altı sekans diyagramı, bir de iki parçalı "
    "sınıf diyagramı çizdik. Şimdi bunlara tek tek bakalım.",
    # 7 Class domain
    "Bu, domain katmanı sınıf diyagramımız — veritabanı varlıklarımız ve ilişkileri. Sağda özetlediğimiz gibi: "
    "bir klinikte birden fazla hasta ve personel var, bu bir-çok ilişki. Hasta ile randevu arasında kompozisyon "
    "var; hasta silinirse randevuları da siliniyor. Her hastanın en fazla bir anamnez kaydı oluyor. Bir randevuda "
    "birden fazla tedavi yapılabiliyor, her tedaviye de bir ödeme bağlanıyor.",
    # 8 Class app
    "Bu da uygulama katmanı. Solda frontend arayüz bileşenleri, yani Boundary sınıfları; ortada hepsinin "
    "haberleştiği ApiService; sağda ise backend kontrolcüleri, ViewSet ve Serializer'lar. Önemli tasarım kararımız "
    "şu: arayüzler doğrudan veritabanıyla konuşmuyor; araya iş mantığını ve güvenliği sağlayan kontrol katmanını "
    "koyduk. Bu da MVC mimarisine tam uyum sağlıyor.",
    # 9 SD-02
    "İlk sekans diyagramımız hasta kaydı. Asistan veya hekim yeni hasta bilgilerini giriyor; istek ApiService "
    "üzerinden PatientViewSet'e, oradan serializer'a gidiyor. Eğer anamnez verisi de varsa, hasta ile birlikte "
    "anamnez kaydı da oluşturuluyor. Güncellemede ise anamnez yoksa otomatik oluşturuluyor.",
    # 10 SD-03
    "Bu, en kritik akışımız: çakışma kontrollü randevu. Asistan randevu kaydetmek istediğinde, serializer "
    "'bu hekimin bu tarih ve saatte planlanmış randevusu var mı' diye kontrol ediyor. Ekrandaki alt-fragmenti "
    "görüyorsunuz: çakışma varsa sistem 400 hatası ve uyarı döndürüyor; çakışma yoksa randevu oluşturulup 201 "
    "dönüyor. Yani FR-03 gereksinimini tam olarak bu akış karşılıyor.",
    # 11 SD-04
    "Dashboard akışı. Personel sabah paneli açtığında, tek bir istekle bugünün randevuları, tamamlanan işlem "
    "sayısı ve toplam hasta sayısı çekiliyor. DashboardView bu verileri filtreleyip tek pakette arayüze dönüyor; "
    "yani gereksiz birden fazla istek atmıyoruz, bu da performans için önemli.",
    # 12 SD-06
    "Son sekans diyagramı klinik ayarları, burada yetkilendirme öne çıkıyor. Yönetici çalışma saatlerini "
    "değiştirmek istediğinde sistem önce yetki kontrolü yapıyor. Kullanıcı admin değilse 403 Forbidden dönüyor; "
    "admin ise ayarlar kaydediliyor ve 200 OK alınıyor. Bu da rol bazlı erişimin somut bir örneği.",
    # 13 DAR backend
    "Teknoloji seçimlerimizi öznel değil, karar matrisleriyle objektif yaptık. Backend için Python ve Django "
    "REST'i seçtik. Tabloda görüldüğü gibi Django REST 19 puanla Express ve Spring'i geçti. En büyük sebep "
    "yerleşik güvenlik ve hazır JWT katmanı — Express'te bunları sıfırdan yazmamız gerekirdi — bir de ekibimizin "
    "Python'a hakimiyeti.",
    # 14 DAR db+frontend
    "Veritabanında PostgreSQL 20 puanla kazandı; çünkü sağlık ve finans verisi ACID uyumu ve kesin ilişkiler "
    "istiyor, MongoDB gibi şemasız sistemler bu yüzden elendi. Frontend'de ise React 20 puan aldı; özellikle "
    "interaktif diş şemamızın DOM performansı ve ekip tecrübesi belirleyici oldu. Sonuçta yığınımız netleşti: "
    "Django REST, React ve PostgreSQL.",
    # 15 API
    "Tüm API'lerimizi Swagger ile canlı dokümante ettik. Birkaç örnek kontrat: giriş uç noktası kullanıcı adı ve "
    "şifre alıp JWT token çifti dönüyor. Randevu uç noktası çakışmada 400, yetkisiz erişimde 403 dönüyor. "
    "Dashboard tek istekte günlük özeti veriyor. Hasta uç noktası ise kayıt ve listeleme yapıyor, eksik alanda 400 dönüyor.",
    # 16 multi-tenant + sms
    "Çok kiracılı yapımız django-tenants ile kurulu. HeaderTenantMiddleware her istekte hangi kliniğe ait "
    "olduğunu okuyup PostgreSQL şemasını ona yönlendiriyor; yani izolasyon fiziksel düzeyde. Dış servis olarak "
    "SMS entegrasyonumuz var: randevu oluşunca hastaya bildirim gidiyor. Dış servis yanıt vermezse işlem çökmüyor, "
    "hata loglanıp randevu asenkron şekilde tamamlanıyor.",
    # 17 test env + critical
    "Test tarafında pytest ve pytest-django kullandık; factory-boy ile Türkçe sahte veri ürettik, monkeypatch ile "
    "de bağımlılıkları taklit ederek metotları gerçekten izole ettik. Sınıf diyagramından, iş mantığı en yoğun olan "
    "5 kritik sınıfı seçtik. Her metot için iki tür senaryo tasarladık: pozitif ve negatif. Önemli nokta şu: negatif "
    "senaryoda sistemin hatayı doğru şekilde reddetmesi başarı sayılıyor, yani durum PASSED.",
    # 18 test matrix
    "İşte sonuç matrisimiz. Her sınıf için temsilî pozitif ve negatif senaryolar var — çakışan randevunun reddi, "
    "eksik alanda 400 dönmesi, mükerrer tedavinin engellenmesi gibi. Hepsi PASSED. Toplamda 25 temsilî senaryonun "
    "tamamı, ilgili üç test dosyasının tam koşumunda ise 54 testin 54'ü başarıyla geçti.",
    # 19 bakım
    "Bakım ayağında üç başlığımız var. Refactoring olarak iş mantığını ayrı bir servis katmanına taşımayı ve "
    "merkezi hata yönetimini planlıyoruz. Yeni özellikler arasında SMS hatırlatma, takvimde sürükle-bırak, "
    "raporlama modülü ve mobil uygulama var. Sağda da v1.1'den v2.1'e kadar sürüm yol haritamızı görüyorsunuz. "
    "Modüler mimari sayesinde bakım maliyetini düşük tutuyoruz.",
    # 20 canlı public
    "Şimdi projenin gerçekten çalıştığını gösterelim — uygulama şu anda canlıda yayında. Bunlar public "
    "sayfalarımız: tanıtım sayfası, klinik kayıt formu ve kliniğe özel giriş ekranı. Yeni bir klinik buradan "
    "kaydolup kendi alt-adresinden sisteme giriş yapabiliyor.",
    # 21 panel 1
    "Giriş yaptıktan sonra karşımıza dashboard çıkıyor: bugünün randevuları, bekleyen ve toplam hasta sayısı tek "
    "ekranda. Yan tarafta hasta yönetimi — arama yapıp hasta listesini görüyoruz. Gördüğünüz gibi tasarımdaki "
    "Boundary sınıfları gerçek arayüze birebir dönüşmüş durumda.",
    # 22 panel 2
    "Burada haftalık randevu takvimimiz var; saat saat planlama yapılabiliyor. Sağda ise çok sekmeli hasta profili: "
    "profil bilgileri, anamnez, tedavi geçmişi, ödemeler, dokümanlar ve en sevdiğimiz kısım — dijital diş şeması, "
    "yani odontogram. Sekans diyagramlarında tasarladığımız akışların çalışan hâli işte bunlar.",
    # 23 sonuç
    "Özetle, Yaşca projesinde yazılım yaşam döngüsünün tüm adımlarını uçtan uca tamamladık: analiz ve UML "
    "modelleme, karar matrisleriyle teknoloji seçimi, API'lerin gerçeklenmesi ve dokümantasyonu, ve 54'te 54 "
    "başarılı birim test. Multi-tenant SaaS mimarisiyle ölçeklenebilir, test edilmiş ve dokümante bir sistem "
    "ortaya koyduk. Dinlediğiniz için teşekkür ederiz; sorularınızı almaktan memnuniyet duyarız.",
]
for _sl, _note in zip(prs.slides, NOTES):
    _sl.notes_slide.notes_text_frame.text = _note

# ---------------------------------------------------------------- kaydet
desktop = os.path.join(os.path.expanduser("~"), "Desktop")
if not os.path.isdir(desktop):
    desktop = os.getcwd()
out = os.path.join(desktop, "Yasca_Final_Proje_Raporu.pptx")
try:
    prs.save(out)
except PermissionError:
    import time
    out = os.path.join(desktop, "Yasca_Final_Proje_Raporu_" + time.strftime("%H%M%S") + ".pptx")
    prs.save(out)
    print("UYARI: Ana dosya acik/kilitli; yeni dosyaya kaydedildi.")
print("KAYDEDILDI:", out)
print("Slayt sayisi:", len(prs.slides._sldIdLst))
